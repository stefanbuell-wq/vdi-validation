#!/usr/bin/env python3
"""
Erstellt eine VDI-Vergleichspräsentation auf Basis des AOK PowerPoint Masters.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy
import os

TEMPLATE = os.path.join(os.path.dirname(__file__), "..", "Präsentation Master.pptx")
OUTPUT = os.path.join(os.path.dirname(__file__), "..", "docs", "VDI-Vergleichsanalyse.pptx")

# AOK brand colors (extracted from template)
AOK_GREEN = RGBColor(0x00, 0x7B, 0x3E)      # AOK Grün
AOK_DARK_GREEN = RGBColor(0x00, 0x5C, 0x2E)  # Dunkelgrün
AOK_LIGHT_GREEN = RGBColor(0xB2, 0xD2, 0x35) # Hellgrün/Pastell
AOK_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AOK_BLACK = RGBColor(0x00, 0x00, 0x00)
AOK_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
AOK_GRAY = RGBColor(0x66, 0x66, 0x66)
AOK_LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
AOK_RED = RGBColor(0xC0, 0x00, 0x00)
AOK_ORANGE = RGBColor(0xE0, 0x7C, 0x00)

# Layout indices from the template
LY_TITLE_GREEN = 2       # 'Titelfolie, grüner Hintergrund'
LY_TITLE_WHITE = 3       # 'Titelfolie, weißer Hintergrund'
LY_SECTION_GREEN = 9     # 'Trennfolie, grüner Hintergrund'
LY_SECTION_WHITE = 10    # 'Trennfolie, weißer Hintergrund'
LY_AGENDA = 6            # 'Agenda'
LY_CONTENT = 12          # 'Titel, Untertitel, Inhalt groß'
LY_TWO_COL = 16          # 'Titel, Untertitel, zwei Inhalte'
LY_THREE_COL = 17        # 'Titel, Untertitel, drei Inhalte'
LY_ONLY_TITLE = 27       # 'Nur Titel'
LY_BLANK = 36            # 'Leeres'


def set_font(run, size=14, bold=False, color=AOK_DARK_GRAY, name="AOK Bundesweit"):
    """Set font properties on a run."""
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    run.font.name = name


def add_paragraph(tf, text, size=14, bold=False, color=AOK_DARK_GRAY, space_after=Pt(6),
                  space_before=Pt(0), alignment=PP_ALIGN.LEFT, level=0, bullet=False):
    """Add a paragraph to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    p.level = level
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    return p


def add_table(slide, rows, cols, left, top, width, height):
    """Add a table shape to a slide."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table


def style_table_cell(cell, text, size=10, bold=False, color=AOK_DARK_GRAY,
                     fill_color=None, alignment=PP_ALIGN.LEFT):
    """Style a table cell."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size=size, bold=bold, color=color)
    cell.text_frame.word_wrap = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color


def style_header_row(table, col_texts, fill_color=AOK_GREEN, text_color=AOK_WHITE):
    """Style the header row of a table."""
    for i, text in enumerate(col_texts):
        style_table_cell(table.cell(0, i), text, size=10, bold=True,
                         color=text_color, fill_color=fill_color,
                         alignment=PP_ALIGN.CENTER)


def set_column_widths(table, widths_cm):
    """Set column widths in cm."""
    for i, w in enumerate(widths_cm):
        table.columns[i].width = Cm(w)


def create_presentation():
    prs = Presentation(TEMPLATE)

    # Remove existing slides (the template has 1 example slide)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # ============================================================
    # SLIDE 1: Title Slide
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TITLE_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = ""
            run = ph.text_frame.paragraphs[0].add_run()
            run.text = "VDI-Lösungen"
            set_font(run, size=40, bold=True, color=AOK_WHITE)
            p2 = ph.text_frame.add_paragraph()
            run2 = p2.add_run()
            run2.text = "Umfassende Vergleichsanalyse"
            set_font(run2, size=28, bold=False, color=AOK_WHITE)
        elif ph.placeholder_format.idx == 1:  # Subtitle
            ph.text = ""
            run = ph.text_frame.paragraphs[0].add_run()
            run.text = "Bewertung für öffentlich-rechtliche Einrichtungen"
            set_font(run, size=18, color=AOK_WHITE)
        elif ph.placeholder_format.idx == 17:  # Bottom left text
            ph.text = ""
            run = ph.text_frame.paragraphs[0].add_run()
            run.text = "März 2026"
            set_font(run, size=12, color=AOK_WHITE)

    # ============================================================
    # SLIDE 2: Agenda
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_AGENDA])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # Title
            ph.text = "Agenda"
        elif ph.placeholder_format.idx == 20:  # Content
            tf = ph.text_frame
            tf.clear()
            # (label, text, indent_level, font_size)
            items = [
                ("1.", "Ausgangslage und Zielsetzung", 0, 18),
                ("2.", "Bewertete VDI-Lösungen im Überblick", 0, 18),
                ("3.", "Einzelanalysen der sechs Lösungen", 0, 18),
                ("4.", "Vergleichsmatrix und TCO-Kostenvergleich", 0, 18),
                ("5.", "Datenschutz und IT-Sicherheit", 0, 18),
                ("6.", "Vergaberechtliche Aspekte", 0, 18),
                ("6a.", "Exkurs: 200 % vs. 120 % Kapazität", 1, 16),
                ("7.", "Bewertungsmatrix und Empfehlung", 0, 18),
                ("8.", "Nächste Schritte", 0, 18),
            ]
            for i, (label, item, level, fsize) in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.level = level
                p.space_after = Pt(10)
                run = p.add_run()
                run.text = f"{label}  {item}"
                set_font(run, size=fsize, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 3: Ausgangslage
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Ausgangslage und Zielsetzung"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Warum eine VDI-Evaluierung?"
        elif ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()
            items = [
                ("Ziel:", " Evaluierung geeigneter VDI-Lösungen für ~7.000 Desktops in unserer öffentlich-rechtlichen Einrichtung"),
                ("Anforderungen:", " Datenschutz (DSGVO), IT-Sicherheit (BSI), Vergaberecht (VgV/UVgO)"),
                ("Bewertungskriterien:", " Kosten, Datenschutz, Funktionalität, Support, Vergaberecht, Zukunftssicherheit"),
                ("Sechs Lösungen:", " Von klassisch On-Premises bis Cloud, von proprietär bis Open Source"),
                ("Zielgruppe:", " Entscheidungsträger IT und Geschäftsführung"),
            ]
            for i, (label, desc) in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(12)
                run_b = p.add_run()
                run_b.text = label
                set_font(run_b, size=15, bold=True, color=AOK_GREEN)
                run_t = p.add_run()
                run_t.text = desc
                set_font(run_t, size=15, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 4: Overview of Solutions
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Bewertete VDI-Lösungen im Überblick"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Sechs Lösungen — von On-Premises bis Cloud"

    # Add overview table
    rows_data = [
        ["Lösung", "Typ", "Hersteller", "Lizenzmodell"],
        ["Windows RDS", "On-Premises", "Microsoft", "CAL-basiert"],
        ["Azure Local", "Hybrid", "Microsoft", "Subscription"],
        ["Azure Virtual Desktop", "Cloud", "Microsoft", "Pay-as-you-go"],
        ["Omnissa Horizon", "On-Prem/Hybrid", "Omnissa (KKR)", "Subscription"],
        ["OpenDesktop", "On-Premises", "Community", "Open Source"],
        ["Proxmox + UDS", "On-Premises", "Proxmox/VirtualCable", "OS + kommerziell"],
    ]

    left, top = Cm(1.5), Cm(4.2)
    width, height = Cm(30), Cm(0.8 * len(rows_data))
    table = add_table(slide, len(rows_data), 4, left, top, width, height)
    set_column_widths(table, [7, 4.5, 8, 5.5])
    style_header_row(table, rows_data[0])
    for r in range(1, len(rows_data)):
        bg = RGBColor(0xF0, 0xF7, 0xF0) if r % 2 == 1 else None
        for c in range(4):
            style_table_cell(table.cell(r, c), rows_data[r][c], size=11, fill_color=bg)

    # ============================================================
    # SLIDE 5: VDI-Bereitstellungsmodelle
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "VDI-Bereitstellungsmodelle"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Der Begriff \u201eVDI\u201c umfasst verschiedene Bereitstellungsmodelle"

    # Modelle-Tabelle
    model_data = [
        ["Modell", "Beschreibung", "Beispiel"],
        ["Desktop-basierte VDI",
         "Ein dedizierter Desktop läuft in einer eigenen VM.\nJeder Nutzer erhält seinen eigenen Desktop.",
         "Proxmox + UDS, Omnissa Horizon,\nCitrix Virtual Desktops"],
        ["Serverbasierte VDI\n(Remote-PC)",
         "Ein Desktop läuft in einem Server-OS,\ndas nur für einen einzelnen Benutzer verfügbar ist.",
         "Citrix Remote PC Access,\nphysische Workstations per Remoting"],
        ["Sitzungsbasierte VDI",
         "Mehrere Benutzer teilen sich ein Server-OS —\njeder Nutzer erhält eine eigene Sitzung.",
         "Windows RDS,\nCitrix Virtual Apps"],
    ]

    left, top = Cm(1.5), Cm(4.2)
    width, height = Cm(30), Cm(0.8 * len(model_data))
    table = add_table(slide, len(model_data), 3, left, top, width, height)
    set_column_widths(table, [7, 13, 10])
    style_header_row(table, model_data[0])
    for r in range(1, len(model_data)):
        bg = RGBColor(0xF0, 0xF7, 0xF0) if r % 2 == 1 else None
        for c in range(3):
            style_table_cell(table.cell(r, c), model_data[r][c], size=10, fill_color=bg)

    # Hinweis unterhalb der Tabelle
    txbox = slide.shapes.add_textbox(Cm(1.5), Cm(8.5), Cm(30), Cm(4.5))
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = ("Wichtig: Bei der aktuellen Infrastruktur werden ca. 1.400 Worker Server "
                "für sitzungsbasierte Desktops betrieben. Bei einer Umstellung auf desktop-basierte VDI "
                "könnten alternativ bis zu 8.000 individuelle Desktops bereitgestellt werden.")
    set_font(run, size=11, bold=False, color=AOK_GRAY)

    add_paragraph(tf, "", size=6, space_after=Pt(4))
    add_paragraph(tf,
        "Citrix als VDI-Plattform: Die bestehende Citrix-Umgebung kann nicht nur sitzungsbasierte "
        "Desktops (Virtual Apps), sondern auch vollwertige Desktop-VDI (Virtual Desktops) bereitstellen. "
        "Ein Wechsel der VDI-Lösung bedeutet nicht zwingend einen Wechsel des Bereitstellungsmodells.",
        size=11, bold=False, color=AOK_GRAY)

    # ============================================================
    # SLIDE 6: Section Divider — Einzelanalysen
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Einzelanalysen"
        elif ph.placeholder_format.idx == 11:
            ph.text = "Detailbewertung der sechs VDI-Lösungen"

    # ============================================================
    # SLIDES 6-11: Individual Solution Analysis (one per solution)
    # ============================================================
    solutions = [
        {
            "name": "Windows Terminal Server (RDS)",
            "subtitle": "Klassische sitzungsbasierte Desktop-Bereitstellung",
            "pros": [
                "Bewährt, breit dokumentiert, großes Ökosystem",
                "Nahtlose AD- und M365-Integration",
                "RDP-Clients auf allen Plattformen",
                "RemoteApps: Einzelanwendungen ohne Full-Desktop",
                "Vorhandenes Know-how in den meisten IT-Abteilungen",
                "Volle On-Premises-Kontrolle",
            ],
            "cons": [
                "Keine echte VDI (sitzungsbasiert, geteiltes OS)",
                "Skalierungsgrenzen (~50-80 User/Server)",
                "Eingeschränkte Multimedia/GPU-Performance",
                "Steigende CAL-Kosten bei Wachstum",
                "Kein modernes zentrales Management",
            ],
            "cost_label": "TCO 5J / 7.000 User",
            "cost_value": "~9,8 Mio. €",
            "cost_per_user": "~280 €/User/Jahr",
        },
        {
            "name": "Azure Local (ehem. Azure Stack HCI)",
            "subtitle": "Hybrid: Azure-Dienste auf eigener Hardware",
            "pros": [
                "Daten bleiben im eigenen RZ",
                "AVD on-premises möglich",
                "Azure-Management (Arc) und -Dienste lokal nutzbar",
                "Skalierbar (2-16 Knoten pro Cluster)",
                "Volle Microsoft-Ökosystem-Kompatibilität",
            ],
            "cons": [
                "Azure-Verbindung für Management zwingend erforderlich",
                "Hohe Einstiegskosten (zertifizierte HCI-Hardware)",
                "Komplexe Lizenzierung (Hardware + Cloud + User)",
                "Starker Microsoft-Lock-in",
                "Azure-Subscription auch bei reinem On-Prem nötig",
            ],
            "cost_label": "TCO 5J / 7.000 User",
            "cost_value": "~28,5 Mio. €",
            "cost_per_user": "~814 €/User/Jahr",
        },
        {
            "name": "Azure Virtual Desktop (Cloud)",
            "subtitle": "Vollständig cloudbasierte VDI in Microsoft Azure",
            "pros": [
                "Keine eigene Hardware nötig",
                "Elastische Skalierung (Autoscaling)",
                "Multi-Session Windows 11 (einzigartig)",
                "Integrierte Sicherheit (Defender, MFA)",
                "FSLogix-Profile, M365-Optimierung",
                "Schnelle Bereitstellung in Minuten",
            ],
            "cons": [
                "Datensouveränität kritisch (CLOUD Act, Schrems-II)",
                "Laufende Kosten bei Dauerbetrieb hoch",
                "Vollständige Internetabhängigkeit",
                "Starker Vendor-Lock-in (Azure)",
                "Komplexe Kostenprognose (verbrauchsbasiert)",
                "Rechtliche Unsicherheit für öffentl. Einrichtungen",
            ],
            "cost_label": "TCO 5J / 7.000 User",
            "cost_value": "~29,8 Mio. €",
            "cost_per_user": "~851 €/User/Jahr",
        },
        {
            "name": "Omnissa Horizon (ehem. VMware)",
            "subtitle": "Enterprise-VDI — seit 07/2024 unter KKR (ex-Broadcom/VMware)",
            "pros": [
                "Vollwertige VDI mit dedizierter VM pro Nutzer",
                "Blast Extreme: Bestes Display-Protokoll (WAN)",
                "Instant Clones, App Volumes, DEM",
                "Windows + Linux als Gast-OS",
                "Unter Omnissa: eigene Roadmap, besserer Support",
                "Flexible Deployment (On-Prem, SaaS, DaaS)",
            ],
            "cons": [
                "vSphere-Kosten (Broadcom): 3-12x teurer geworden",
                "Nur noch Subscription (keine Perpetual-Lizenzen)",
                "Komplexe Infrastruktur (vSphere, vCenter, UAG)",
                "Unsichere Langfrist-Strategie (KKR = Private Equity)",
                "Basic Support abgekündigt (nur noch Production)",
                "Starker Vendor-Lock-in",
            ],
            "cost_label": "TCO 5J / 7.000 User",
            "cost_value": "~34,0 Mio. €",
            "cost_per_user": "~971 €/User/Jahr",
        },
        {
            "name": "Open Source — OpenDesktop",
            "subtitle": "KVM/QEMU + Guacamole + XRDP + oVirt",
            "pros": [
                "Keine Lizenzkosten (Kernkomponenten frei)",
                "Maximale Transparenz (Quellcode einsehbar)",
                "Volle Datensouveränität, keine Telemetrie",
                "Vendor-unabhängig, fördert Wettbewerb",
                "Linux-native, BSI-Audit möglich",
                "Vergaberechtlich vorteilhaft",
            ],
            "cons": [
                "Kein einheitliches Produkt — Integration nötig",
                "Windows als Gast erfordert weiterhin MS-Lizenzen",
                "Kein kommerzieller Support aus einer Hand",
                "Hoher Integrationsaufwand und Personalkosten",
                "Display-Performance hinter proprietären Protokollen",
                "Wenig Referenzen in deutschen Behörden",
            ],
            "cost_label": "TCO 5J / 7.000 User",
            "cost_value": "~16,5 Mio. €",
            "cost_per_user": "~471 €/User/Jahr",
        },
        {
            "name": "Proxmox VE + UDS Enterprise",
            "subtitle": "Open-Source-Hypervisor + professioneller VDI-Broker (EU-Hersteller)",
            "pros": [
                "Proxmox kostenlos nutzbar (Community Edition)",
                "UDS Enterprise: Enterprise-VDI-Features",
                "EU-Hersteller (Österreich/Spanien) — DSGVO-konform",
                "Multi-Protokoll (RDP, SPICE, HTML5, X2Go)",
                "Kosteneffizient vs. VMware/Citrix",
                "Referenzen im öffentlichen Sektor (Spanien)",
                "Enterprise-Support über Fujitsu als bestehender Dienstleister verfügbar",
            ],
            "cons": [
                "Geringere Marktpräsenz in Deutschland",
                "UDS Enterprise ist kommerziell (nicht Open Source)",
                "Kein BSI-/CC-Zertifikat für Proxmox",
                "Kleineres Ökosystem und weniger Partner",
                "Kein FSLogix/App Volumes-Äquivalent",
                "Know-how-Aufbau in DACH noch begrenzt",
            ],
            "cost_label": "TCO 5J / 7.000 User",
            "cost_value": "~15,5 Mio. €",
            "cost_per_user": "~443 €/User/Jahr",
        },
    ]

    for sol in solutions:
        slide = prs.slides.add_slide(prs.slide_layouts[LY_TWO_COL])
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:  # Title
                ph.text = sol["name"]
            elif ph.placeholder_format.idx == 15:  # Subtitle
                ph.text = sol["subtitle"]
            elif ph.placeholder_format.idx == 16:  # Left content - Vorteile
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "Vorteile"
                set_font(run, size=14, bold=True, color=AOK_GREEN)
                p.space_after = Pt(6)
                for item in sol["pros"]:
                    p = tf.add_paragraph()
                    p.level = 0
                    p.space_after = Pt(3)
                    run = p.add_run()
                    run.text = f"+ {item}"
                    set_font(run, size=10, color=AOK_DARK_GRAY)
                # Cost info at bottom
                p = tf.add_paragraph()
                p.space_before = Pt(12)
                run = p.add_run()
                run.text = f"{sol['cost_label']}: {sol['cost_value']} ({sol['cost_per_user']})"
                set_font(run, size=11, bold=True, color=AOK_GREEN)

            elif ph.placeholder_format.idx == 17:  # Right content - Nachteile
                tf = ph.text_frame
                tf.clear()
                p = tf.paragraphs[0]
                run = p.add_run()
                run.text = "Nachteile"
                set_font(run, size=14, bold=True, color=AOK_RED)
                p.space_after = Pt(6)
                for item in sol["cons"]:
                    p = tf.add_paragraph()
                    p.level = 0
                    p.space_after = Pt(3)
                    run = p.add_run()
                    run.text = f"– {item}"
                    set_font(run, size=10, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 12: Section Divider — Vergleich
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Vergleich und\nKostenanalyse"
        elif ph.placeholder_format.idx == 11:
            ph.text = "Gegenüberstellung aller sechs Lösungen"

    # ============================================================
    # SLIDE 13: Comparison Matrix
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Vergleichsmatrix — Kerneigenschaften"
        elif ph.placeholder_format.idx == 15:
            ph.text = ""

    headers = ["Kriterium", "RDS", "Azure\nLocal", "AVD\n(Cloud)", "Omnissa\nHorizon", "Open\nDesktop", "Proxmox\n+ UDS"]
    data = [
        ["Deployment", "On-Prem", "Hybrid", "Cloud", "On-Prem/\nHybrid", "On-Prem", "On-Prem"],
        ["Echte VDI", "Nein", "Ja", "Ja", "Ja", "Ja", "Ja"],
        ["Multi-Session", "Ja", "Ja", "Ja", "Ja", "Eingeschr.", "Eingeschr."],
        ["Linux-Desktop", "Eingeschr.", "Ja", "Eingeschr.", "Ja", "Ja (nativ)", "Ja (nativ)"],
        ["Display-Protokoll", "RDP", "RDP", "RDP opt.", "Blast\nExtreme", "VNC/SPICE", "SPICE/RDP\n/HTML5"],
        ["WAN-Perform.", "Mittel", "Mittel-Gut", "Gut", "Sehr gut", "Gering-\nMittel", "Mittel"],
        ["Automatisierung", "Gering", "Hoch", "Sehr hoch", "Hoch", "Gering", "Mittel"],
        ["Vendor-Lock-in", "Mittel", "Hoch", "Sehr hoch", "Hoch", "Keiner", "Gering"],
        ["Datensouveränität", "Hoch", "Mittel-\nHoch", "Gering", "Hoch", "Sehr hoch", "Sehr hoch"],
        ["Reifegrad", "Sehr hoch", "Hoch", "Hoch", "Sehr hoch", "Mittel", "Mittel-\nHoch"],
    ]

    n_rows = len(data) + 1
    n_cols = 7
    left, top = Cm(1.5), Cm(3.5)
    width, height = Cm(30.5), Cm(0.75 * n_rows)
    table = add_table(slide, n_rows, n_cols, left, top, width, height)
    set_column_widths(table, [4.5, 3.5, 3.5, 3.5, 4, 3.5, 4])
    style_header_row(table, headers)

    # Color-code cells based on content
    color_map = {
        "Ja": RGBColor(0xE8, 0xF5, 0xE9),
        "Ja (nativ)": RGBColor(0xE8, 0xF5, 0xE9),
        "Sehr hoch": RGBColor(0xE8, 0xF5, 0xE9),
        "Hoch": RGBColor(0xF1, 0xF8, 0xE9),
        "Sehr gut": RGBColor(0xE8, 0xF5, 0xE9),
        "Gut": RGBColor(0xF1, 0xF8, 0xE9),
        "Nein": RGBColor(0xFF, 0xEB, 0xEE),
        "Gering": RGBColor(0xFF, 0xEB, 0xEE),
        "Eingeschr.": RGBColor(0xFF, 0xF3, 0xE0),
        "Keiner": RGBColor(0xE8, 0xF5, 0xE9),
        "Sehr hoch": RGBColor(0xE8, 0xF5, 0xE9),
    }

    for r, row_data in enumerate(data):
        for c, cell_text in enumerate(row_data):
            clean = cell_text.replace("\n", "")
            bg = color_map.get(clean, None) if c > 0 else RGBColor(0xF5, 0xF5, 0xF5)
            bold = c == 0
            align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            style_table_cell(table.cell(r + 1, c), cell_text, size=9, bold=bold,
                             fill_color=bg, alignment=align)

    # ============================================================
    # SLIDE 14: TCO Cost Comparison
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "TCO-Kostenvergleich — 5 Jahre, 7.000 Benutzer"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Total Cost of Ownership inkl. Hardware, Lizenzen, Support, Personal (in Mio. €)"

    tco_headers = ["Kostenposition", "RDS", "Azure\nLocal", "AVD\n(Cloud)", "Omnissa\nHorizon", "Open\nDesktop", "Proxmox\n+ UDS"]
    tco_data = [
        ["Hardware (Server,\nStorage, Netzwerk)", "2,0 Mio.", "5,0 Mio.", "0 €", "4,5 Mio.", "3,5 Mio.", "3,5 Mio."],
        ["Lizenzen (5 J.)", "3,5 Mio.", "16,0 Mio.", "18,0 Mio.", "18,0 Mio.", "0 €", "2,5 Mio."],
        ["Windows-Lizenzen (5 J.)", "inkl. EA", "inkl.", "inkl.", "3,5 Mio.", "3,5 Mio.", "3,5 Mio."],
        ["Support (5 J.)", "0,8 Mio.", "2,0 Mio.", "1,0 Mio.", "2,5 Mio.", "2,5 Mio.", "1,5 Mio."],
        ["Implementierung", "0,5 Mio.", "1,5 Mio.", "0,8 Mio.", "1,5 Mio.", "2,0 Mio.", "1,0 Mio."],
        ["Internes Personal\n(5 J., ~6-10 FTE)", "3,0 Mio.", "4,0 Mio.", "2,5 Mio.", "4,0 Mio.", "5,0 Mio.", "3,5 Mio."],
        ["TCO 5 Jahre", "~9,8 Mio.", "~28,5 Mio.", "~29,8 Mio.", "~34,0 Mio.", "~16,5 Mio.", "~15,5 Mio."],
        ["TCO pro User/Jahr", "~280 €", "~814 €", "~851 €", "~971 €", "~471 €", "~443 €"],
    ]

    n_rows = len(tco_data) + 1
    left, top = Cm(1.5), Cm(3.5)
    width, height = Cm(30.5), Cm(0.7 * n_rows)
    table = add_table(slide, n_rows, 7, left, top, width, height)
    set_column_widths(table, [5.5, 3.5, 3.5, 3.5, 3.5, 3.5, 4])
    style_header_row(table, tco_headers)

    for r, row_data in enumerate(tco_data):
        is_total = r >= len(tco_data) - 2
        for c, cell_text in enumerate(row_data):
            bg = AOK_GREEN if is_total else (RGBColor(0xF0, 0xF7, 0xF0) if r % 2 == 0 else None)
            txt_color = AOK_WHITE if is_total else AOK_DARK_GRAY
            bold = c == 0 or is_total
            align = PP_ALIGN.RIGHT if c > 0 else PP_ALIGN.LEFT
            style_table_cell(table.cell(r + 1, c), cell_text, size=9, bold=bold,
                             color=txt_color, fill_color=bg, alignment=align)

    # ============================================================
    # SLIDE 15: Zusätzliche Hardware-Investitionen (PoC, Migration, GPU)
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Zusätzliche Hardware-Investitionen"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Separat zu kalkulieren: Test-/PoC-Hardware, Migrations-Hardware, GPU-Ausstattung"

    # Three-column table: PoC | Migration | GPU
    hw_headers = ["Kostenblock", "Beschreibung", "Geschätzte Kosten", "Status"]
    hw_data = [
        ["Test-/PoC-Hardware",
         "Dedizierte Server (2–4 Hosts/Lösung),\nStorage, Netzwerk, Test-Endgeräte\nfür PoC mit 2–3 Lösungen",
         "65.000–190.000 €\n(einmalig, in Produktion\nüberführbar)",
         "Zu beschaffen"],
        ["Migrations-Hardware\n(Parallelbetrieb)",
         "Zusätzliche Kapazität während\nMigration (Alt- + Neusystem parallel)\nBei phasenweiser Migration: 15–30 %",
         "300.000–600.000 €\n(phasenweise, 7.000 User)\nBig Bang: 1,5–3,0 Mio. €",
         "Menge noch zu\nbestimmen"],
        ["GPU-Hardware\n(falls benötigt)",
         "vGPU-Karten (z.B. NVIDIA A16/L4)\n+ GPU-fähige Server-Hosts\n+ NVIDIA vGPU-Lizenzen",
         "Bei 10 % GPU-Nutzern:\n~280.000–468.000 € (5 J.)\nOhne GPU: 0 €",
         "Bedarfsanalyse\nerforderlich"],
    ]

    n_rows = len(hw_data) + 1
    left, top = Cm(1.5), Cm(3.5)
    width, height = Cm(30.5), Cm(1.3 * n_rows)
    table = add_table(slide, n_rows, 4, left, top, width, height)
    set_column_widths(table, [5.5, 9, 8, 5])
    style_header_row(table, hw_headers)

    status_colors = {
        "Zu beschaffen": AOK_ORANGE,
        "Menge noch zu\nbestimmen": AOK_RED,
        "Bedarfsanalyse\nerforderlich": AOK_RED,
    }

    for r, row_data in enumerate(hw_data):
        bg = RGBColor(0xF0, 0xF7, 0xF0) if r % 2 == 0 else None
        for c, cell_text in enumerate(row_data):
            bold = c == 0
            align = PP_ALIGN.LEFT
            txt_color = AOK_DARK_GRAY
            if c == 3:  # Status column
                txt_color = status_colors.get(cell_text, AOK_DARK_GRAY)
                bold = True
                align = PP_ALIGN.CENTER
            elif c == 2:  # Cost column
                align = PP_ALIGN.RIGHT
                bold = True
            style_table_cell(table.cell(r + 1, c), cell_text, size=9, bold=bold,
                             color=txt_color, fill_color=bg, alignment=align)

    # Add note below table
    from pptx.util import Cm as CmUtil
    txBox = slide.shapes.add_textbox(Cm(1.5), Cm(3.5 + 1.3 * n_rows + 0.5), Cm(30.5), Cm(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Offene Fragen:"
    set_font(run, size=11, bold=True, color=AOK_RED)

    questions = [
        "Wird GPU-Power benötigt? → Benutzerprofilanalyse durchführen (CAD, GIS, PACS, Multimedia, KI)",
        "Welche Migrationsstrategie? → Bestimmt Menge der Migrations-Hardware (Big Bang vs. phasenweise)",
        "Kann PoC-Hardware in Produktion übernommen werden? → Bei Beschaffung auf Prod-Eignung achten",
    ]
    for q in questions:
        p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = f"• {q}"
        set_font(run, size=10, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 16: Section Divider — Datenschutz
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Datenschutz und\nIT-Sicherheit"
        elif ph.placeholder_format.idx == 11:
            ph.text = "DSGVO, BSI, Schrems-II, CLOUD Act"

    # ============================================================
    # SLIDE 17: Data Sovereignty Assessment
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Datensouveränität — Risikobewertung"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Schrems-II / CLOUD Act Analyse für öffentlich-rechtliche Einrichtungen"

    ds_headers = ["Lösung", "Datenstandort", "US-Zugriff\nmöglich?", "Schrems-II\nRisiko", "Bewertung"]
    ds_data = [
        ["RDS", "Eigenes RZ", "Nein", "Gering", "Sehr gut"],
        ["Azure Local", "Eigenes RZ\n(Mgmt: Azure)", "Metadaten:\nJa", "Mittel", "Gut"],
        ["AVD (Cloud)", "Azure-RZ\n(EU möglich)", "Ja\n(CLOUD Act)", "Hoch", "Kritisch"],
        ["Omnissa Horizon", "Eigenes RZ", "Nein\n(On-Prem)", "Gering", "Sehr gut"],
        ["OpenDesktop", "Eigenes RZ", "Nein", "Sehr gering", "Sehr gut"],
        ["Proxmox + UDS", "Eigenes RZ", "Nein", "Sehr gering", "Sehr gut"],
    ]

    n_rows = len(ds_data) + 1
    left, top = Cm(1.5), Cm(3.5)
    width, height = Cm(30.5), Cm(0.9 * n_rows)
    table = add_table(slide, n_rows, 5, left, top, width, height)
    set_column_widths(table, [5.5, 5, 5, 5, 5])
    style_header_row(table, ds_headers)

    risk_colors = {
        "Sehr gut": RGBColor(0xE8, 0xF5, 0xE9),
        "Gut": RGBColor(0xF1, 0xF8, 0xE9),
        "Kritisch": RGBColor(0xFF, 0xEB, 0xEE),
        "Gering": RGBColor(0xE8, 0xF5, 0xE9),
        "Sehr gering": RGBColor(0xE8, 0xF5, 0xE9),
        "Mittel": RGBColor(0xFF, 0xF3, 0xE0),
        "Hoch": RGBColor(0xFF, 0xEB, 0xEE),
    }

    for r, row_data in enumerate(ds_data):
        for c, cell_text in enumerate(row_data):
            bg = None
            if c == 3:  # Risiko column
                bg = risk_colors.get(cell_text.strip(), None)
            elif c == 4:  # Bewertung column
                bg = risk_colors.get(cell_text.strip(), None)
            bold = c == 0
            align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            style_table_cell(table.cell(r + 1, c), cell_text, size=11, bold=bold,
                             fill_color=bg, alignment=align)

    # ============================================================
    # SLIDE 18: Datenschutz Empfehlungen
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Datenschutz — Empfehlungen und Rahmenbedingungen"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Regulatorischer Rahmen für öffentlich-rechtliche Einrichtungen"
        elif ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()
            sections = [
                ("Bevorzugt:", " On-Premises-Lösungen (RDS, Omnissa Horizon, Proxmox + UDS, OpenDesktop)"),
                ("Bedingt geeignet:", " Azure Local (Hybrid) — nach DSFA und mit organisatorischen Maßnahmen"),
                ("Kritisch:", " AVD Cloud — nur nach umfassender DSFA und Genehmigung durch DSB"),
                ("", ""),
                ("BSI C5:2025:", " Erweiterte Kontrollen, Souveränitätskriterien — verpflichtend ab 01/2027"),
                ("Delos Cloud:", " Souveräne Azure-Alternative unter deutscher Betriebsführung (SAP/Delos GmbH)"),
                ("DVC:", " Deutsche Verwaltungscloud seit 04/2025 produktiv — vereinfachte Beschaffung"),
                ("", ""),
                ("Landesdatenschutzbeauftragte:", " Einsatz von MS 365/Azure in öffentl. Verwaltung als problematisch eingestuft"),
            ]
            for i, (label, desc) in enumerate(sections):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                if not label and not desc:
                    p.space_after = Pt(4)
                    continue
                p.space_after = Pt(6)
                if label:
                    run_b = p.add_run()
                    run_b.text = label
                    color = AOK_GREEN if "Bevorzugt" in label or "BSI" in label or "Delos" in label or "DVC" in label else (AOK_ORANGE if "Bedingt" in label else AOK_RED)
                    set_font(run_b, size=12, bold=True, color=color)
                run_t = p.add_run()
                run_t.text = desc
                set_font(run_t, size=12, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 19: Section Divider — Vergaberecht
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Vergaberechtliche\nAspekte"
        elif ph.placeholder_format.idx == 11:
            ph.text = "EU-Schwellenwerte, VgV, UVgO, EVB-IT"

    # ============================================================
    # SLIDE 20: Vergaberecht
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TWO_COL])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Vergaberecht — Schwellenwerte und Verfahren"
        elif ph.placeholder_format.idx == 15:
            ph.text = "EU-Schwellenwerte seit 01.01.2026"
        elif ph.placeholder_format.idx == 16:  # Left
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "EU-Schwellenwerte (netto)"
            set_font(run, size=13, bold=True, color=AOK_GREEN)

            items = [
                ("Bundesbehörden:", " 140.000 €"),
                ("Sonstige öffentl. AG:", " 216.000 €"),
                ("Sektorenauftraggeber:", " 432.000 €"),
                ("", ""),
                ("Unterschwellenbereich (UVgO):", ""),
                ("Direktvergabe:", " bis 100.000 € möglich"),
                ("", ""),
                ("Praxis:", " Bei 7.000 Desktops (TCO 10-34 Mio. €) ist eine EU-weite Ausschreibung zwingend erforderlich"),
            ]
            for label, val in items:
                p = tf.add_paragraph()
                p.space_after = Pt(4)
                if not label and not val:
                    continue
                if label:
                    run = p.add_run()
                    run.text = label
                    set_font(run, size=11, bold=True, color=AOK_DARK_GRAY)
                if val:
                    run = p.add_run()
                    run.text = val
                    set_font(run, size=11, color=AOK_DARK_GRAY)

        elif ph.placeholder_format.idx == 17:  # Right
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "EVB-IT Vertragstypen"
            set_font(run, size=13, bold=True, color=AOK_GREEN)

            evb_items = [
                "EVB-IT System — On-Premises VDI + Hardware",
                "EVB-IT Cloud — Cloud-VDI (AVD, Azure Local); C5-Testat erforderlich",
                "EVB-IT Überlassung Typ B — Softwarelizenzen (z.B. UDS Enterprise)",
                "EVB-IT Pflege S — Wartung und Patching",
                "EVB-IT Dienstvertrag — Implementierung/Beratung",
                "EVB-IT Rahmenvereinbarung — Modulare VDI-Leistungen",
                "",
                "Pflicht: EVB-IT für Bundesbehörden (§ 55 BHO)",
            ]
            for item in evb_items:
                p = tf.add_paragraph()
                p.space_after = Pt(4)
                if not item:
                    continue
                run = p.add_run()
                run.text = item
                bold = "Pflicht" in item
                set_font(run, size=10, bold=bold, color=AOK_RED if bold else AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 21: Ausschreibungshinweise
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Ausschreibungshinweise"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Empfohlene Zuschlagskriterien und produktneutrale Vergabe"
        elif ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()
            items = [
                ("Produktneutrale Ausschreibung:", " Funktionale Anforderungen definieren, kein Produkt vorschreiben (§ 31 Abs. 6 VgV)"),
                ("Losaufteilung:", " Hardware, Software, Dienstleistungen ggf. in separate Lose (§ 97 Abs. 4 GWB)"),
                ("", ""),
                ("Empfohlene Zuschlagskriterien:", ""),
                ("  Preis:", " 30–40 %"),
                ("  Technische Leistungsfähigkeit:", " 20–30 %"),
                ("  Datenschutz / IT-Sicherheit:", " 15–25 %"),
                ("  Support und SLA:", " 10–15 %"),
                ("  Migrationsfähigkeit / Exit-Strategie:", " 5–10 %"),
            ]
            for i, (label, val) in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(6)
                if not label and not val:
                    continue
                if label:
                    run = p.add_run()
                    run.text = label
                    set_font(run, size=13, bold=True, color=AOK_GREEN if "Empfohlene" not in label else AOK_GREEN)
                if val:
                    run = p.add_run()
                    run.text = val
                    set_font(run, size=13, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 22: Section Divider — Exkurs Kapazität
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Exkurs:\n200 % vs. 120 %"
        elif ph.placeholder_format.idx == 11:
            ph.text = "Brauchen wir wirklich die volle Verdopplung?"

    # ============================================================
    # SLIDE 23: 200 % vs. 120 % — Modellvergleich
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "200 % vs. 120 % — Kapazitätsmodelle im Vergleich"

    # Model comparison table
    cap_headers = ["Kriterium", "200 % (Ist-Zustand)", "120 % (Alternative)"]
    cap_data = [
        ["Infrastruktur-Faktor", "2,0× der Basiskapazität", "1,2× der Basiskapazität"],
        ["Prinzip", "100 % pro RZ\n(Active-Passive)", "100 % Betrieb + 20 % Puffer\n(Active-Active)"],
        ["Hardwarekosten", "2× Basis", "1,2× Basis"],
        ["Lizenzkosten", "2× Basis", "1,2× Basis"],
        ["Strom / RZ-Fläche", "2× Basis", "1,2× Basis"],
        ["Verfügbarkeit bei\nRZ-Ausfall", "100 % — voller Betrieb\nohne Einschränkung", "Eingeschränkt —\nPriorisierung nötig"],
        ["Betankung /\nProvisioning", "Jederzeit volle Kapazität", "20 % Puffer deckt\nlaufende Rollouts ab"],
        ["Performance-Puffer", "Großzügig", "20 % Headroom —\nnormale Lastspitzen"],
        ["Kostenersparnis", "Referenz", "ca. 40 % weniger"],
    ]

    rows = len(cap_data) + 1
    cols = len(cap_headers)
    table = add_table(slide, rows, cols, Cm(1.5), Cm(3.5), Cm(22), Cm(13))
    set_column_widths(table, [5.5, 8.25, 8.25])
    style_header_row(table, cap_headers)
    for r, row_data in enumerate(cap_data, start=1):
        for c, val in enumerate(row_data):
            bold = (c == 0)
            bg = RGBColor(0xF0, 0xF8, 0xF0) if c == 2 else None
            style_table_cell(table.cell(r, c), val, size=9, bold=bold,
                             color=AOK_GREEN if c == 2 and r == len(cap_data) else AOK_DARK_GRAY,
                             fill_color=bg)

    # ============================================================
    # SLIDE 24: 200 % vs. 120 % — Kostenbeispiel und Risiko
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_TWO_COL])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "200 % vs. 120 % — Kosten und Risikobewertung"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Rechenbeispiel und Mitigationsmaßnahmen"
        elif ph.placeholder_format.idx == 14:
            # Left column: cost example
            tf = ph.text_frame
            tf.clear()
            add_paragraph(tf, "Rechenbeispiel (Proxmox + UDS, 5 Jahre)", size=13, bold=True, color=AOK_GREEN, space_after=Pt(10))
            cost_lines = [
                ("Hardware:", "200 %: ~200 T€  →  120 %: ~120 T€"),
                ("Lizenzen & Support:", "200 %: ~250 T€  →  120 %: ~150 T€"),
                ("Strom & RZ-Fläche:", "200 %: ~50 T€  →  120 %: ~30 T€"),
                ("Gesamt:", "200 %: ~500 T€  →  120 %: ~300 T€"),
            ]
            for label, val in cost_lines:
                p = tf.add_paragraph()
                p.space_after = Pt(4)
                run = p.add_run()
                run.text = label + "  "
                is_total = "Gesamt" in label
                set_font(run, size=11, bold=True, color=AOK_GREEN if is_total else AOK_DARK_GRAY)
                run = p.add_run()
                run.text = val
                set_font(run, size=11, bold=is_total, color=AOK_GREEN if is_total else AOK_DARK_GRAY)

            p = tf.add_paragraph()
            p.space_after = Pt(10)
            add_paragraph(tf, "→ Potenzielle Ersparnis: ca. 200.000 € (40 %)", size=12, bold=True, color=AOK_GREEN, space_after=Pt(6))

        elif ph.placeholder_format.idx == 16:
            # Right column: risk assessment
            tf = ph.text_frame
            tf.clear()
            add_paragraph(tf, "Risiko bei 120 % — Mitigation", size=13, bold=True, color=AOK_GREEN, space_after=Pt(10))
            risk_items = [
                ("Priorisierte Wiederherstellung:", " Tier-1-Desktops (geschäftskritisch) zuerst, Tier-2 zeitversetzt"),
                ("Elastische Skalierung:", " Proxmox HA verteilt VMs in Minuten auf verbleibende Hosts"),
                ("Degraded Mode:", " Kurzzeitig reduzierte Ressourcen pro VM akzeptabel"),
            ]
            for label, val in risk_items:
                p = tf.add_paragraph()
                p.space_after = Pt(8)
                run = p.add_run()
                run.text = label
                set_font(run, size=11, bold=True, color=AOK_DARK_GRAY)
                run = p.add_run()
                run.text = val
                set_font(run, size=11, color=AOK_DARK_GRAY)

            p = tf.add_paragraph()
            p.space_after = Pt(12)
            add_paragraph(tf, "Empfehlung", size=13, bold=True, color=AOK_GREEN, space_after=Pt(6))
            add_paragraph(tf, "Workshop mit RZ-Betrieb, IT-Sicherheit und Einkauf — SLA-Anforderungen definieren, reale Lastdaten analysieren, Kapazitätspuffer bestimmen.", size=11, color=AOK_DARK_GRAY, space_after=Pt(4))

    # ============================================================
    # SLIDE 25: Section Divider — Empfehlung
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_SECTION_GREEN])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Bewertung und\nEmpfehlung"
        elif ph.placeholder_format.idx == 11:
            ph.text = "Gewichtete Bewertungsmatrix"

    # ============================================================
    # SLIDE 23: Scoring Matrix
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_ONLY_TITLE])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Gewichtete Bewertungsmatrix"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Skala: 1 (schlecht) bis 10 (sehr gut)"

    score_headers = ["Kriterium (Gewicht)", "RDS", "Azure\nLocal", "AVD\n(Cloud)", "Omnissa\nHorizon", "Open\nDesktop", "Proxmox\n+ UDS"]
    score_data = [
        ["Kosten (25 %)", "9", "5", "5", "3", "7", "8"],
        ["Datenschutz (25 %)", "8", "6", "3", "7", "10", "9"],
        ["Funktionsumfang (15 %)", "4", "8", "9", "9", "4", "7"],
        ["Support (10 %)", "8", "7", "8", "5", "3", "6"],
        ["Vergaberecht (10 %)", "7", "6", "5", "4", "9", "8"],
        ["Zukunftssicherheit (10 %)", "6", "8", "8", "4", "6", "7"],
        ["Betriebsaufwand (5 %)", "6", "5", "9", "5", "2", "6"],
        ["Gewichtete Summe", "7,25", "6,30", "5,90", "5,40", "6,75", "7,70"],
    ]

    n_rows = len(score_data) + 1
    left, top = Cm(1.5), Cm(3.5)
    width, height = Cm(30.5), Cm(0.85 * n_rows)
    table = add_table(slide, n_rows, 7, left, top, width, height)
    set_column_widths(table, [6, 3.5, 3.5, 3.5, 3.5, 3.5, 4])
    style_header_row(table, score_headers)

    def score_color(val_str):
        try:
            v = float(val_str.replace(",", "."))
            if v >= 8:
                return RGBColor(0xC8, 0xE6, 0xC9)
            elif v >= 6:
                return RGBColor(0xE8, 0xF5, 0xE9)
            elif v >= 4:
                return RGBColor(0xFF, 0xF3, 0xE0)
            else:
                return RGBColor(0xFF, 0xEB, 0xEE)
        except ValueError:
            return None

    for r, row_data in enumerate(score_data):
        is_total = r == len(score_data) - 1
        for c, cell_text in enumerate(row_data):
            if is_total:
                bg = AOK_GREEN
                txt_color = AOK_WHITE
            elif c == 0:
                bg = RGBColor(0xF5, 0xF5, 0xF5)
                txt_color = AOK_DARK_GRAY
            else:
                bg = score_color(cell_text)
                txt_color = AOK_DARK_GRAY
            bold = c == 0 or is_total
            align = PP_ALIGN.CENTER if c > 0 else PP_ALIGN.LEFT
            style_table_cell(table.cell(r + 1, c), cell_text, size=11, bold=bold,
                             color=txt_color, fill_color=bg, alignment=align)

    # ============================================================
    # SLIDE 24: Ranking & Recommendation
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Rangfolge und Handlungsempfehlung"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Ergebnis der gewichteten Bewertung"
        elif ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()

            rankings = [
                ("1.", "Proxmox + UDS Enterprise", "7,70", "Bestes Gesamtpaket: Kosten, Datenschutz, Skalierbarkeit — TCO ~15,5 Mio. €"),
                ("2.", "Windows RDS", "7,25", "Günstigste Lösung (~9,8 Mio. €), aber sitzungsbasiert — keine echte VDI"),
                ("3.", "OpenDesktop", "6,75", "Maximale Datensouveränität, aber ~10 FTE Personalaufwand"),
                ("4.", "Azure Local", "6,30", "Guter Kompromiss, aber ~28,5 Mio. € und starker MS-Lock-in"),
                ("5.", "Azure Virtual Desktop", "5,90", "Autoscaling-Vorteil bei 7.000 Usern, Datenschutz kritisch"),
                ("6.", "Omnissa Horizon", "5,40", "Technisch ausgereift, aber teuerste Option (~34 Mio. €)"),
            ]

            for i, (rank, name, score, note) in enumerate(rankings):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(8)
                # Rank
                run = p.add_run()
                run.text = f"{rank} "
                color = AOK_GREEN if i < 2 else (AOK_ORANGE if i < 4 else AOK_RED)
                set_font(run, size=14, bold=True, color=color)
                # Name + Score
                run = p.add_run()
                run.text = f"{name} ({score})  "
                set_font(run, size=14, bold=True, color=AOK_DARK_GRAY)
                # Note
                run = p.add_run()
                run.text = f"— {note}"
                set_font(run, size=12, color=AOK_GRAY)

    # ============================================================
    # SLIDE 25: Handlungsempfehlung
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Handlungsempfehlung"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Strategische Empfehlung für die VDI-Einführung"
        elif ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()

            recommendations = [
                ("Kurzfristig / einfache Anforderungen:",
                 "Windows RDS als bewährte Lösung — TCO ~9,8 Mio. € für 7.000 User, aber keine echte VDI und begrenzte Skalierungsfähigkeit (~100 RDS-Hosts nötig)"),
                ("Empfehlung für 7.000 Desktops:",
                 "Proxmox VE + UDS Enterprise — TCO ~15,5 Mio. €, echte VDI, datenschutzkonform, EU-Hersteller; bei dieser Größenordnung der beste Kosten-Nutzen-Kompromiss"),
                ("Strategisch / Digitale Souveränität:",
                 "OpenDesktop-Ansatz als langfristige Open-Source-Strategie (~16,5 Mio. €), aber ~10 FTE nötig — nur mit entsprechendem Personalaufbau realistisch"),
                ("Mit Vorsicht:",
                 "Omnissa Horizon (~34 Mio. €) und AVD Cloud (~29,8 Mio. €) — erhebliches Kostenrisiko bei 7.000 Desktops; AVD zusätzlich Datenschutzrisiko (CLOUD Act)"),
            ]

            for i, (title, desc) in enumerate(recommendations):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(4)
                p.space_before = Pt(8)
                run = p.add_run()
                run.text = title
                color = AOK_GREEN if i < 3 else AOK_ORANGE
                set_font(run, size=14, bold=True, color=color)

                p2 = tf.add_paragraph()
                p2.space_after = Pt(8)
                run2 = p2.add_run()
                run2.text = desc
                set_font(run2, size=13, color=AOK_DARK_GRAY)

    # ============================================================
    # SLIDE 26: Next Steps
    # ============================================================
    slide = prs.slides.add_slide(prs.slide_layouts[LY_CONTENT])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Nächste Schritte"
        elif ph.placeholder_format.idx == 15:
            ph.text = "Fahrplan zur VDI-Einführung"
        elif ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()

            steps = [
                ("1.", "Anforderungsanalyse", "7.000 Desktops: Nutzerprofile, Anwendungslandschaft, Performance-Klassen definieren"),
                ("2.", "Benutzerprofilanalyse GPU", "Klären: Wie viele Nutzer benötigen GPU-Power? (CAD, GIS, PACS, Multimedia, KI)"),
                ("3.", "Test-/PoC-Hardware beschaffen", "Dedizierte Hardware für Proof of Concept (ca. 65.000–190.000 €)"),
                ("4.", "Proof of Concept (PoC)", "Pilotierung mit 2–3 favorisierten Lösungen"),
                ("5.", "Migrationsstrategie festlegen", "Parallelbetrieb-Kapazität und Migrations-Hardware bestimmen"),
                ("6.", "DSFA + Abstimmung mit DSB", "Datenschutz-Folgenabschätzung und Datenschutzbeauftragten einbeziehen"),
                ("7.", "Marktrecherche + EU-Ausschreibung", "Markterkundung (§ 28 VgV), dann EU-weite Vergabe (TCO >> Schwellenwert)"),
            ]

            for i, (num, title, desc) in enumerate(steps):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.space_after = Pt(2)
                run = p.add_run()
                run.text = f"{num}  {title}"
                set_font(run, size=14, bold=True, color=AOK_GREEN)

                p2 = tf.add_paragraph()
                p2.space_after = Pt(8)
                run2 = p2.add_run()
                run2.text = f"     {desc}"
                set_font(run2, size=12, color=AOK_GRAY)

    # ============================================================
    # Save
    # ============================================================
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Presentation saved to: {OUTPUT}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    create_presentation()
