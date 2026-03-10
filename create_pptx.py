#!/usr/bin/env python3
"""
Erstellt STACKIT VDI-Kostenanalyse Präsentation auf Basis des AOK Präsentation Master.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
import math

# AOK Brand Colors (from master)
AOK_GREEN = RGBColor(0x00, 0x50, 0x32)       # Dunkelgrün
AOK_LIGHT_GREEN = RGBColor(0x00, 0xA0, 0x5E) # Hellgrün
AOK_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AOK_BLACK = RGBColor(0x00, 0x00, 0x00)
AOK_GRAY = RGBColor(0x80, 0x80, 0x80)
AOK_DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
HIGHLIGHT_GREEN = RGBColor(0x00, 0xB0, 0x50)
HIGHLIGHT_RED = RGBColor(0xC0, 0x00, 0x00)

prs = Presentation("Präsentation Master.pptx")

# Remove the existing template slide
if len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId:
        prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# ----- Helper Functions -----

def add_table(slide, rows, cols, left, top, width, height):
    """Add a table to a slide."""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    return table_shape.table

def set_cell(table, row, col, text, font_size=10, bold=False, alignment=PP_ALIGN.LEFT,
             font_color=AOK_BLACK, fill_color=None):
    """Set cell text with formatting."""
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = font_color
    run.font.name = "AOK Bureau Office"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Reduce margins
    cell.margin_left = Cm(0.15)
    cell.margin_right = Cm(0.15)
    cell.margin_top = Cm(0.05)
    cell.margin_bottom = Cm(0.05)
    if fill_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_color

def set_header_row(table, texts, font_size=9):
    """Format header row with green background."""
    for i, text in enumerate(texts):
        set_cell(table, 0, i, text, font_size=font_size, bold=True,
                font_color=AOK_WHITE, fill_color=AOK_GREEN, alignment=PP_ALIGN.CENTER)

def add_content_slide(prs, title_text, layout_idx=12):
    """Add a content slide using 'Titel, Untertitel, Inhalt groß' layout."""
    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)
    # Set title
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            shape.text = title_text
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(22)
                    run.font.bold = True
            break
    return slide

def add_text_to_placeholder(slide, ph_idx, text, font_size=14, bold=False, color=AOK_BLACK):
    """Add text to a specific placeholder."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            shape.text = ""
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.color.rgb = color
            return shape
    return None

def add_bullet_text(slide, ph_idx, items, font_size=12, bold_first=False):
    """Add bulleted text to a placeholder."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == ph_idx:
            tf = shape.text_frame
            tf.clear()
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = item
                run.font.size = Pt(font_size)
                run.font.color.rgb = AOK_BLACK
                if bold_first and i == 0:
                    run.font.bold = True
                p.space_after = Pt(4)
            return shape
    return None

def add_textbox(slide, left, top, width, height, text, font_size=11,
                bold=False, color=AOK_BLACK, alignment=PP_ALIGN.LEFT):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "AOK Bureau Office"
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=11, color=AOK_BLACK):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, size_override) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size_override if size_override else font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "AOK Bureau Office"
        p.space_after = Pt(3)
    return txBox


# ============================================================
# SLIDE 1: Titelfolie
# ============================================================
layout = prs.slide_layouts[2]  # Titelfolie, grüner Hintergrund
slide = prs.slides.add_slide(layout)
for shape in slide.placeholders:
    idx = shape.placeholder_format.idx
    if idx == 0:  # Titel
        shape.text = "STACKIT VDI-Kostenanalyse"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(32)
                r.font.bold = True
    elif idx == 1:  # Untertitel
        shape.text = "7.000 Arbeitsplätze mit Linux Mint & openDesk"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(18)
    elif idx == 17:  # Textplatzhalter unten
        shape.text = "Auf Basis der STACKIT Preisliste v1.0.36 (Stand 03.03.2026)"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(11)

# ============================================================
# SLIDE 2: Agenda
# ============================================================
layout = prs.slide_layouts[6]  # Agenda
slide = prs.slides.add_slide(layout)
for shape in slide.placeholders:
    idx = shape.placeholder_format.idx
    if idx == 0:
        shape.text = "Agenda"
    elif idx == 20:
        tf = shape.text_frame
        tf.clear()
        agenda_items = [
            "1.  Ausgangslage & Annahmen",
            "2.  VM-Sizing: Drei Szenarien",
            "3.  Compute-Kosten (On-Demand vs. 24/7)",
            "4.  Storage-Kosten",
            "5.  Infrastruktur & Netzwerk",
            "6.  Gesamtkosten-Übersicht",
            "7.  Vergleich: Citrix on-premise vs. STACKIT",
            "8.  Fazit & Empfehlung",
        ]
        for i, item in enumerate(agenda_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = item
            run.font.size = Pt(16)
            run.font.color.rgb = AOK_BLACK
            run.font.name = "AOK Bureau Office"
            p.space_after = Pt(8)

# ============================================================
# SLIDE 3: Ausgangslage & Annahmen
# ============================================================
slide = add_content_slide(prs, "Ausgangslage & Annahmen", layout_idx=27)

# Annahmen-Tabelle
tbl = add_table(slide, 9, 3,
    Cm(1.5), Cm(2.5), Cm(22), Cm(10))

set_header_row(tbl, ["Parameter", "Wert", "Begründung"], font_size=10)

data = [
    ("Gesamtnutzer", "7.000", "Vorgabe"),
    ("Concurrent-User-Rate", "70%", "Urlaub, Teilzeit, Meetings, Krankheit"),
    ("Gleichzeitige Sessions", "4.900", "7.000 x 0,70"),
    ("Betriebssystem", "Linux Mint", "Keine Windows-Lizenz erforderlich"),
    ("Arbeitszeit/Tag", "10 Stunden", "07:00 - 17:00 (mit Puffer)"),
    ("Arbeitstage/Monat", "22 Tage", "Standard"),
    ("Nutzungsstunden/Monat", "220 h", "10h x 22 Tage (On-Demand)"),
    ("STACKIT Referenzmonat", "720 h", "30-Tage-Monat Kalkulationsbasis"),
]

col_widths = [Cm(6), Cm(4), Cm(12)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for i, (param, wert, grund) in enumerate(data):
    row = i + 1
    bg = RGBColor(0xF0, 0xF8, 0xF0) if row % 2 == 0 else None
    set_cell(tbl, row, 0, param, font_size=10, bold=True, fill_color=bg)
    set_cell(tbl, row, 1, wert, font_size=10, bold=True, alignment=PP_ALIGN.CENTER,
             font_color=AOK_GREEN, fill_color=bg)
    set_cell(tbl, row, 2, grund, font_size=10, fill_color=bg)

# ============================================================
# SLIDE 4: VM-Sizing Szenarien
# ============================================================
slide = add_content_slide(prs, "VM-Sizing: Drei Szenarien", layout_idx=27)

tbl = add_table(slide, 4, 5,
    Cm(1.5), Cm(2.5), Cm(22), Cm(5.5))

set_header_row(tbl, ["Szenario", "vCPU", "RAM", "Typische Nutzung", "STACKIT SKU"], font_size=10)

col_widths = [Cm(4), Cm(2), Cm(2.5), Cm(8.5), Cm(5)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

scenarios = [
    ("A: Basis-Worker", "2", "8 GB",
     "E-Mail, Browser, openDesk, SAP GUI for Java",
     "g1.2-EU01 (Intel)"),
    ("B: Standard-Worker", "4", "16 GB",
     "Wie A + mehrere parallele Anwendungen, leichte Auswertungen",
     "g1.3-EU01 (Intel)"),
    ("C: Power-User", "8", "32 GB",
     "Wie B + BI-Tools, Datenverarbeitung, Entwicklungsumgebungen",
     "g1.4-EU01 (Intel)"),
]

colors = [
    RGBColor(0xE8, 0xF5, 0xE9),
    RGBColor(0xFF, 0xF3, 0xE0),
    RGBColor(0xFC, 0xE4, 0xEC),
]

for i, (sz, vcpu, ram, usage, sku) in enumerate(scenarios):
    row = i + 1
    set_cell(tbl, row, 0, sz, font_size=10, bold=True, fill_color=colors[i])
    set_cell(tbl, row, 1, vcpu, font_size=10, alignment=PP_ALIGN.CENTER, fill_color=colors[i])
    set_cell(tbl, row, 2, ram, font_size=10, alignment=PP_ALIGN.CENTER, fill_color=colors[i])
    set_cell(tbl, row, 3, usage, font_size=9, fill_color=colors[i])
    set_cell(tbl, row, 4, sku, font_size=9, bold=True, fill_color=colors[i])

# Zusatzinfo
add_textbox(slide, Cm(1.5), Cm(8.5), Cm(22), Cm(3),
    "Empfehlung: Intel g1-Generation ist ca. 35% günstiger als AMD g2a.\n"
    "Für VDI-Workloads ohne spezielle AMD-Anforderungen ist g1 die beste Wahl.\n\n"
    "Hinweis: '-m' Suffix = mit Windows-Lizenz (hier nicht benötigt wegen Linux Mint).",
    font_size=10, color=AOK_GRAY)

# ============================================================
# SLIDE 5: Compute-Kosten Vergleich
# ============================================================
slide = add_content_slide(prs, "Compute-Kosten: On-Demand vs. 24/7", layout_idx=27)

# On-Demand Tabelle
add_textbox(slide, Cm(1.5), Cm(2.3), Cm(12), Cm(0.8),
    "On-Demand Betrieb (220h/Monat) - EMPFOHLEN", font_size=12, bold=True, color=AOK_GREEN)

tbl = add_table(slide, 4, 5,
    Cm(1.5), Cm(3.2), Cm(22), Cm(4.5))

set_header_row(tbl, ["Szenario", "Preis/h", "x 220h", "x 4.900 VMs", "Gesamt/Monat"], font_size=9)

col_widths = [Cm(5), Cm(3.5), Cm(3.5), Cm(4), Cm(6)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

ondemand_data = [
    ("A: Basis (2/8)", "0,0758 EUR", "16,68 EUR", "4.900", "81.718 EUR"),
    ("B: Standard (4/16)", "0,1516 EUR", "33,36 EUR", "4.900", "163.464 EUR"),
    ("C: Power (8/32)", "0,3033 EUR", "66,73 EUR", "4.900", "326.957 EUR"),
]

for i, (sz, ph, m, vms, total) in enumerate(ondemand_data):
    row = i + 1
    set_cell(tbl, row, 0, sz, font_size=10, bold=True)
    set_cell(tbl, row, 1, ph, font_size=10, alignment=PP_ALIGN.RIGHT)
    set_cell(tbl, row, 2, m, font_size=10, alignment=PP_ALIGN.RIGHT)
    set_cell(tbl, row, 3, vms, font_size=10, alignment=PP_ALIGN.CENTER)
    set_cell(tbl, row, 4, total, font_size=10, bold=True, alignment=PP_ALIGN.RIGHT,
             font_color=AOK_GREEN)

# 24/7 Tabelle
add_textbox(slide, Cm(1.5), Cm(8.0), Cm(12), Cm(0.8),
    "24/7 Betrieb (720h/Monat)", font_size=12, bold=True, color=AOK_GRAY)

tbl2 = add_table(slide, 4, 4,
    Cm(1.5), Cm(8.9), Cm(22), Cm(4))

set_header_row(tbl2, ["Szenario", "Preis/Monat (720h)", "x 4.900 VMs", "Gesamt/Monat"], font_size=9)

col_widths2 = [Cm(5), Cm(5.5), Cm(4), Cm(7.5)]
for i, w in enumerate(col_widths2):
    tbl2.columns[i].width = w

fulltime_data = [
    ("A: Basis (2/8)", "54,59 EUR", "4.900", "267.491 EUR"),
    ("B: Standard (4/16)", "109,18 EUR", "4.900", "534.982 EUR"),
    ("C: Power (8/32)", "218,37 EUR", "4.900", "1.070.013 EUR"),
]

for i, (sz, pm, vms, total) in enumerate(fulltime_data):
    row = i + 1
    set_cell(tbl2, row, 0, sz, font_size=10, bold=True)
    set_cell(tbl2, row, 1, pm, font_size=10, alignment=PP_ALIGN.RIGHT)
    set_cell(tbl2, row, 2, vms, font_size=10, alignment=PP_ALIGN.CENTER)
    set_cell(tbl2, row, 3, total, font_size=10, alignment=PP_ALIGN.RIGHT, font_color=AOK_GRAY)

# Einsparung Highlight
add_textbox(slide, Cm(14), Cm(13.2), Cm(10), Cm(0.8),
    "Einsparung durch On-Demand: ca. 69%",
    font_size=12, bold=True, color=AOK_GREEN, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 6: Storage-Kosten
# ============================================================
slide = add_content_slide(prs, "Storage-Kosten", layout_idx=27)

add_textbox(slide, Cm(1.5), Cm(2.3), Cm(12), Cm(0.8),
    "Block Storage pro VDI-Desktop (persistent für alle 7.000 User)", font_size=11, bold=True, color=AOK_GREEN)

tbl = add_table(slide, 5, 5,
    Cm(1.5), Cm(3.2), Cm(22), Cm(5))

set_header_row(tbl, ["Komponente", "Menge", "SKU / Typ", "Preis/Einheit/Monat", "Gesamt/Monat"], font_size=9)

col_widths = [Cm(5), Cm(3.5), Cm(5.5), Cm(4), Cm(4)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

storage_data = [
    ("OS + Apps Disk", "50 GB/User", "Premium-Capacity", "0,07 EUR/GB", "24.500 EUR"),
    ("Performance Tier", "1 Disk/User", "Perf. 2 (1000 IOPS)", "14,50 EUR/Disk", "101.500 EUR"),
    ("Object Storage", "10 GB/User = 70 TB", "Object Storage Premium", "0,03 EUR/GB", "2.100 EUR"),
    ("File Storage (Profile)", "5 GB/User = 35 TB", "File Storage Standard", "0,22 EUR/GB", "7.700 EUR"),
]

for i, (comp, menge, sku, preis, total) in enumerate(storage_data):
    row = i + 1
    bg = RGBColor(0xF0, 0xF8, 0xF0) if row % 2 == 0 else None
    set_cell(tbl, row, 0, comp, font_size=9, bold=True, fill_color=bg)
    set_cell(tbl, row, 1, menge, font_size=9, alignment=PP_ALIGN.CENTER, fill_color=bg)
    set_cell(tbl, row, 2, sku, font_size=9, fill_color=bg)
    set_cell(tbl, row, 3, preis, font_size=9, alignment=PP_ALIGN.RIGHT, fill_color=bg)
    set_cell(tbl, row, 4, total, font_size=9, bold=True, alignment=PP_ALIGN.RIGHT, fill_color=bg)

# Total
add_textbox(slide, Cm(14), Cm(8.5), Cm(9.5), Cm(0.8),
    "Storage Gesamt: 135.800 EUR/Monat",
    font_size=13, bold=True, color=AOK_GREEN, alignment=PP_ALIGN.RIGHT)

add_textbox(slide, Cm(1.5), Cm(9.5), Cm(22), Cm(1.5),
    "Hinweis: Storage wird für alle 7.000 User berechnet (persistente Daten),\n"
    "nicht nur für die 4.900 concurrent Sessions.",
    font_size=10, color=AOK_GRAY)


# ============================================================
# SLIDE 7: Infrastruktur & Netzwerk
# ============================================================
slide = add_content_slide(prs, "Infrastruktur- & Netzwerk-Kosten", layout_idx=27)

add_textbox(slide, Cm(1.5), Cm(2.3), Cm(12), Cm(0.8),
    "Management-Server (24/7 Betrieb)", font_size=11, bold=True, color=AOK_GREEN)

tbl = add_table(slide, 7, 5,
    Cm(1.5), Cm(3.2), Cm(22), Cm(6))

set_header_row(tbl, ["Rolle", "Anzahl", "SKU (vCPU/RAM)", "Preis/VM/Monat", "Gesamt/Monat"], font_size=9)

col_widths = [Cm(6.5), Cm(2), Cm(5.5), Cm(4), Cm(4)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

infra_data = [
    ("VDI Connection Broker", "3", "g1.3 (4/16)", "109,18 EUR", "328 EUR"),
    ("openDesk Stack", "5", "g1.4 (8/32)", "218,37 EUR", "1.092 EUR"),
    ("Identity Management", "2", "g1.3 (4/16)", "109,18 EUR", "218 EUR"),
    ("Monitoring / Logging", "2", "g1.3 (4/16)", "109,18 EUR", "218 EUR"),
    ("Reverse Proxy / Gateway", "2", "g1.2 (2/8)", "54,59 EUR", "109 EUR"),
    ("GESAMT (14 VMs)", "", "", "", "1.965 EUR"),
]

for i, (rolle, anz, sku, preis, total) in enumerate(infra_data):
    row = i + 1
    is_total = (i == len(infra_data) - 1)
    bg = AOK_GREEN if is_total else (RGBColor(0xF0, 0xF8, 0xF0) if row % 2 == 0 else None)
    fc = AOK_WHITE if is_total else AOK_BLACK
    set_cell(tbl, row, 0, rolle, font_size=9, bold=True, fill_color=bg, font_color=fc)
    set_cell(tbl, row, 1, anz, font_size=9, alignment=PP_ALIGN.CENTER, fill_color=bg, font_color=fc)
    set_cell(tbl, row, 2, sku, font_size=9, fill_color=bg, font_color=fc)
    set_cell(tbl, row, 3, preis, font_size=9, alignment=PP_ALIGN.RIGHT, fill_color=bg, font_color=fc)
    set_cell(tbl, row, 4, total, font_size=9, bold=True, alignment=PP_ALIGN.RIGHT, fill_color=bg, font_color=fc)

# Netzwerk
add_textbox(slide, Cm(1.5), Cm(9.8), Cm(12), Cm(0.8),
    "Netzwerk", font_size=11, bold=True, color=AOK_GREEN)

tbl2 = add_table(slide, 4, 4,
    Cm(1.5), Cm(10.7), Cm(22), Cm(3.5))

set_header_row(tbl2, ["Komponente", "Anzahl", "Preis/Monat", "Gesamt/Monat"], font_size=9)

col_w2 = [Cm(8), Cm(4), Cm(5), Cm(5)]
for i, w in enumerate(col_w2):
    tbl2.columns[i].width = w

nw_data = [
    ("Public IP (IPv4)", "10", "2,92 EUR", "29 EUR"),
    ("App Load Balancer (10k conn)", "4", "18,92 EUR", "76 EUR"),
    ("DNS Zone (1000 Records)", "2", "3,09 EUR", "6 EUR"),
]

for i, (comp, anz, preis, total) in enumerate(nw_data):
    row = i + 1
    set_cell(tbl2, row, 0, comp, font_size=9)
    set_cell(tbl2, row, 1, anz, font_size=9, alignment=PP_ALIGN.CENTER)
    set_cell(tbl2, row, 2, preis, font_size=9, alignment=PP_ALIGN.RIGHT)
    set_cell(tbl2, row, 3, total, font_size=9, bold=True, alignment=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 8: Gesamtkosten-Übersicht
# ============================================================
slide = add_content_slide(prs, "Gesamtkosten-Übersicht (On-Demand)", layout_idx=27)

tbl = add_table(slide, 8, 4,
    Cm(1.5), Cm(2.5), Cm(22), Cm(8))

set_header_row(tbl, ["Kostenblock", "A: Basis (2/8)", "B: Standard (4/16)", "C: Power (8/32)"], font_size=10)

col_widths = [Cm(7), Cm(5), Cm(5), Cm(5)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

overview_data = [
    ("Compute (4.900 VMs, 220h)", "81.718 EUR", "163.464 EUR", "326.957 EUR"),
    ("Block Storage (7.000 User)", "126.000 EUR", "126.000 EUR", "126.000 EUR"),
    ("Shared Storage", "9.800 EUR", "9.800 EUR", "9.800 EUR"),
    ("Infrastruktur-Server", "1.965 EUR", "1.965 EUR", "1.965 EUR"),
    ("Netzwerk", "111 EUR", "111 EUR", "111 EUR"),
    ("MONATLICH GESAMT", "~219.594 EUR", "~301.340 EUR", "~464.833 EUR"),
    ("PRO USER / MONAT", "~31,37 EUR", "~43,05 EUR", "~66,40 EUR"),
]

for i, (block, a, b, c) in enumerate(overview_data):
    row = i + 1
    is_total = (i >= len(overview_data) - 2)
    bg = AOK_GREEN if is_total else (RGBColor(0xF0, 0xF8, 0xF0) if row % 2 == 0 else None)
    fc = AOK_WHITE if is_total else AOK_BLACK
    gc = AOK_WHITE if is_total else AOK_GREEN
    set_cell(tbl, row, 0, block, font_size=10, bold=(is_total or i==0), fill_color=bg, font_color=fc)
    set_cell(tbl, row, 1, a, font_size=10, alignment=PP_ALIGN.RIGHT, fill_color=bg, font_color=gc if is_total else AOK_BLACK, bold=is_total)
    set_cell(tbl, row, 2, b, font_size=10, alignment=PP_ALIGN.RIGHT, fill_color=bg, font_color=gc if is_total else AOK_BLACK, bold=is_total)
    set_cell(tbl, row, 3, c, font_size=10, alignment=PP_ALIGN.RIGHT, fill_color=bg, font_color=gc if is_total else AOK_BLACK, bold=is_total)

add_textbox(slide, Cm(1.5), Cm(11), Cm(22), Cm(1),
    "Alle Preise netto zzgl. MwSt. | STACKIT Preisliste v1.0.36 | Mengenrabatte nicht berücksichtigt",
    font_size=9, color=AOK_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: Kosten pro User/Monat Zusammenfassung
# ============================================================
slide = add_content_slide(prs, "Kosten pro User / Monat", layout_idx=27)

tbl = add_table(slide, 4, 3,
    Cm(3), Cm(3), Cm(18), Cm(5))

set_header_row(tbl, ["Szenario", "24/7 Betrieb", "On-Demand (220h)"], font_size=12)

col_widths = [Cm(7), Cm(5.5), Cm(5.5)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

user_cost_data = [
    ("A: Basis (2 vCPU / 8 GB)", "~58 EUR", "~31 EUR"),
    ("B: Standard (4 vCPU / 16 GB)", "~97 EUR", "~43 EUR"),
    ("C: Power (8 vCPU / 32 GB)", "~185 EUR", "~66 EUR"),
]

for i, (sz, full, od) in enumerate(user_cost_data):
    row = i + 1
    set_cell(tbl, row, 0, sz, font_size=12, bold=True)
    set_cell(tbl, row, 1, full, font_size=12, alignment=PP_ALIGN.CENTER, font_color=AOK_GRAY)
    set_cell(tbl, row, 2, od, font_size=14, bold=True, alignment=PP_ALIGN.CENTER,
             font_color=AOK_GREEN, fill_color=RGBColor(0xE8, 0xF5, 0xE9))

add_textbox(slide, Cm(3), Cm(8.5), Cm(18), Cm(1),
    "On-Demand spart ca. 69% gegenüber 24/7 Betrieb",
    font_size=14, bold=True, color=AOK_GREEN, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Cm(3), Cm(10), Cm(18), Cm(2),
    "Jahreskosten (On-Demand):\n"
    "A: ~2,64 Mio. EUR  |  B: ~3,62 Mio. EUR  |  C: ~5,58 Mio. EUR",
    font_size=11, color=AOK_DARK_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: Vergleich Citrix vs. STACKIT
# ============================================================
slide = add_content_slide(prs, "Vergleich: Citrix on-premise vs. STACKIT Linux VDI", layout_idx=27)

tbl = add_table(slide, 8, 3,
    Cm(2), Cm(2.5), Cm(20), Cm(8.5))

set_header_row(tbl, ["Kostenposten", "Citrix (typisch)", "STACKIT Linux VDI"], font_size=10)

col_widths = [Cm(7), Cm(6.5), Cm(6.5)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

compare_data = [
    ("Citrix-Lizenzen", "15-25 EUR/User/Mon", "0 EUR (Open Source VDI)"),
    ("Windows-Lizenzen (VDA)", "7-12 EUR/User/Mon", "0 EUR (Linux Mint)"),
    ("MS Office-Lizenzen", "10-15 EUR/User/Mon", "0 EUR (openDesk/Collabora)"),
    ("Hardware / Datacenter", "15-30 EUR/User/Mon", "In Compute enthalten"),
    ("Compute + Storage (Cloud)", "n/a", "31-66 EUR/User/Mon"),
    ("Netzwerk + Infra", "~3-5 EUR/User/Mon", "~0,30 EUR/User/Mon"),
    ("GESCHÄTZT GESAMT", "~50-80 EUR/User/Mon", "~31-66 EUR/User/Mon"),
]

for i, (posten, citrix, stackit) in enumerate(compare_data):
    row = i + 1
    is_total = (i == len(compare_data) - 1)
    bg = AOK_GREEN if is_total else (RGBColor(0xF0, 0xF8, 0xF0) if row % 2 == 0 else None)
    fc = AOK_WHITE if is_total else AOK_BLACK

    set_cell(tbl, row, 0, posten, font_size=10, bold=True, fill_color=bg, font_color=fc)

    # Citrix column - red-ish for costs
    citrix_color = AOK_WHITE if is_total else HIGHLIGHT_RED
    set_cell(tbl, row, 1, citrix, font_size=10, alignment=PP_ALIGN.CENTER,
             font_color=citrix_color, fill_color=bg, bold=is_total)

    # STACKIT column - green for savings
    stackit_color = AOK_WHITE if is_total else HIGHLIGHT_GREEN
    set_cell(tbl, row, 2, stackit, font_size=10, alignment=PP_ALIGN.CENTER,
             font_color=stackit_color, fill_color=bg, bold=is_total)

add_textbox(slide, Cm(2), Cm(11.5), Cm(20), Cm(1),
    "STACKIT + Linux + openDesk: bis zu 50-60% günstiger als Citrix on-premise",
    font_size=14, bold=True, color=AOK_GREEN, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: Fazit & Empfehlung
# ============================================================
layout = prs.slide_layouts[9]  # Trennfolie, grüner Hintergrund
slide = prs.slides.add_slide(layout)

for shape in slide.placeholders:
    idx = shape.placeholder_format.idx
    if idx == 0:
        shape.text = "Fazit & Empfehlung"
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(28)
                r.font.bold = True
    elif idx == 11:
        tf = shape.text_frame
        tf.clear()
        items = [
            ("Empfohlenes Szenario: B (Standard, 4 vCPU / 16 GB)", True, 16),
            ("", False, 8),
            ("Kostenvorteil:", True, 14),
            ("  ~43 EUR/User/Monat vs. 50-80 EUR bei Citrix", False, 13),
            ("  Keine Lizenzkosten für OS, Office oder VDI-Broker", False, 13),
            ("  On-Demand Betrieb spart zusätzlich 69%", False, 13),
            ("", False, 8),
            ("Nächste Schritte:", True, 14),
            ("  1. PoC mit 50 Usern auf STACKIT aufsetzen", False, 13),
            ("  2. Mengenrabatte mit STACKIT verhandeln", False, 13),
            ("  3. openDesk-Kompatibilität validieren", False, 13),
            ("  4. Migrationspfad von Citrix definieren", False, 13),
        ]
        for i, (text, bold, size) in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            p.space_after = Pt(2)

# ============================================================
# SLIDE 12: Nicht enthalten / Disclaimer
# ============================================================
slide = add_content_slide(prs, "Nicht in dieser Kalkulation enthalten", layout_idx=27)

items_left = [
    "VDI-Broker-Software (z.B. Apache Guacamole = kostenlos)",
    "openDesk-Betriebskosten (Personal, Wartung)",
    "Bandbreite / Egress-Traffic",
    "Backup-Kosten (ca. 0,03 EUR/GB/Monat)",
    "Projektkosten (Migration, Setup, Schulung)",
    "Support-Vertrag STACKIT",
    "Mengenrabatte (bei 7.000 Usern verhandelbar!)",
]

y = Cm(3)
for item in items_left:
    add_textbox(slide, Cm(2.5), y, Cm(20), Cm(0.7),
        "   " + item, font_size=11, color=AOK_DARK_GRAY)
    y += Cm(1.1)

add_textbox(slide, Cm(1.5), Cm(12), Cm(22), Cm(1),
    "Quelle: STACKIT Preisliste v1.0.36 (Stand 03.03.2026) | Alle Preise netto zzgl. MwSt.",
    font_size=9, color=AOK_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# Save
# ============================================================
output_path = "STACKIT_VDI_Kostenanalyse_7000_User.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
